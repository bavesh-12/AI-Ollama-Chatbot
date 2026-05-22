import os
import pickle
import faiss
import numpy as np
import hashlib
import fitz
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
from docx import Document
import pdfplumber
import mimetypes
import chardet
from pptx import Presentation
import markdown
import re
import pandas as pd
from PIL import Image
import pytesseract
import zipfile
import io

VECTOR_DIR = "vector_store"
os.makedirs(VECTOR_DIR, exist_ok=True)

embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
EMBEDDING_DIM = 384

def get_conversation_store_path(conversation_id: str):
    conv_dir = os.path.join(VECTOR_DIR, conversation_id)
    os.makedirs(conv_dir, exist_ok=True)
    return (
        os.path.join(conv_dir, "index.faiss"),
        os.path.join(conv_dir, "data.pkl"),
        os.path.join(conv_dir, "metadata.pkl")
    )

def initialize_vector_store(conversation_id: str):
    INDEX_FILE, DATA_FILE, METADATA_FILE = get_conversation_store_path(conversation_id)
    
    index = faiss.IndexFlatL2(EMBEDDING_DIM)
    documents = []
    metadata = []
    
    faiss.write_index(index, INDEX_FILE)
    with open(DATA_FILE, "wb") as f:
        pickle.dump(documents, f)
    with open(METADATA_FILE, "wb") as f:
        pickle.dump(metadata, f)
    
    return index, documents, metadata

def load_vector_store(conversation_id: str):
    INDEX_FILE, DATA_FILE, METADATA_FILE = get_conversation_store_path(conversation_id)
    
    if not os.path.exists(INDEX_FILE):
        return initialize_vector_store(conversation_id)
    
    try:
        index = faiss.read_index(INDEX_FILE)
        with open(DATA_FILE, "rb") as f:
            documents = pickle.load(f)
        with open(METADATA_FILE, "rb") as f:
            metadata = pickle.load(f)
        return index, documents, metadata
    except:
        return initialize_vector_store(conversation_id)

def get_embedding(text: str):
    return embedding_model.encode(text).astype("float32")

def clean_text(text: str):
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n+', '\n', text)
    text = text.strip()
    return text

def extract_with_pymupdf(file_path: str):
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except:
        pass
    return clean_text(text)

def extract_with_pdfplumber(file_path: str):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except:
        pass
    return clean_text(text)

def extract_with_pypdf2(file_path: str):
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except:
        pass
    return clean_text(text)

def extract_text_from_pdf(file_path: str):
    best_text = ""
    best_length = 0
    
    methods = [
        ("PyMuPDF", extract_with_pymupdf),
        ("PDFPlumber", extract_with_pdfplumber),
        ("PyPDF2", extract_with_pypdf2)
    ]
    
    for method_name, method_func in methods:
        try:
            extracted = method_func(file_path)
            if extracted and len(extracted) > best_length:
                best_length = len(extracted)
                best_text = extracted
        except:
            continue
    
    return best_text

def extract_text_from_docx(file_path: str):
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + "\n"
    except:
        pass
    return clean_text(text)

def extract_text_from_pptx(file_path: str):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
    except:
        pass
    return clean_text(text)

def extract_text_from_txt(file_path: str):
    text = ""
    try:
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] if result['encoding'] else 'utf-8'
        
        with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
            text = f.read()
    except:
        try:
            with open(file_path, 'r', encoding='latin-1', errors='ignore') as f:
                text = f.read()
        except:
            pass
    return clean_text(text)

def extract_text_from_md(file_path: str):
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        text = md_content
    except:
        text = extract_text_from_txt(file_path)
    return clean_text(text)

def extract_text_from_image(file_path: str):
    text = ""
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
    except:
        pass
    return clean_text(text)

def extract_text_from_csv(file_path: str):
    text = ""
    try:
        df = pd.read_csv(file_path)
        text = df.to_string()
    except:
        pass
    return clean_text(text)

def extract_text_from_excel(file_path: str):
    text = ""
    try:
        df = pd.read_excel(file_path)
        text = df.to_string()
    except:
        pass
    return clean_text(text)

def extract_text_from_zip(file_path: str):
    text = ""
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            for file_name in zip_ref.namelist():
                if file_name.endswith(('.txt', '.md', '.csv', '.json')):
                    with zip_ref.open(file_name) as f:
                        content = f.read().decode('utf-8', errors='ignore')
                        text += f"File: {file_name}\n{content}\n\n"
    except:
        pass
    return clean_text(text)

def extract_text_from_file(file_path: str, filename: str):
    text = ""
    
    try:
        if not os.path.exists(file_path):
            return 0, "", "File not found"
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return 0, "", "Empty file"
        
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif file_ext == '.docx':
            text = extract_text_from_docx(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            text = extract_text_from_pptx(file_path)
        elif file_ext == '.txt':
            text = extract_text_from_txt(file_path)
        elif file_ext in ['.md', '.markdown']:
            text = extract_text_from_md(file_path)
        elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff']:
            text = extract_text_from_image(file_path)
        elif file_ext == '.csv':
            text = extract_text_from_csv(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            text = extract_text_from_excel(file_path)
        elif file_ext in ['.zip', '.rar']:
            text = extract_text_from_zip(file_path)
        else:
            mime_type, _ = mimetypes.guess_type(file_path)
            if mime_type and 'pdf' in mime_type:
                text = extract_text_from_pdf(file_path)
            elif mime_type and 'word' in mime_type:
                text = extract_text_from_docx(file_path)
            elif mime_type and 'powerpoint' in mime_type:
                text = extract_text_from_pptx(file_path)
            elif mime_type and 'text' in mime_type:
                text = extract_text_from_txt(file_path)
            elif mime_type and 'image' in mime_type:
                text = extract_text_from_image(file_path)
            else:
                text = extract_text_from_txt(file_path)
        
        if text and len(text.strip()) > 10:
            return 1, "file_hash", text
        else:
            return 0, "", text if text else "Could not extract text"
            
    except Exception as e:
        return 0, "", f"Error: {str(e)}"

def add_to_memory(conversation_id: str, text: str, source: str, source_id: str = "", chunk_num: int = 0):
    try:
        index, documents, metadata = load_vector_store(conversation_id)
        INDEX_FILE, DATA_FILE, METADATA_FILE = get_conversation_store_path(conversation_id)
        
        embedding = get_embedding(text).reshape(1, -1)
        index.add(embedding)
        documents.append(text)
        metadata.append({
            "source": source,
            "source_id": source_id,
            "chunk_num": chunk_num
        })
        
        faiss.write_index(index, INDEX_FILE)
        with open(DATA_FILE, "wb") as f:
            pickle.dump(documents, f)
        with open(METADATA_FILE, "wb") as f:
            pickle.dump(metadata, f)
    except Exception as e:
        pass

def add_document_to_memory(conversation_id: str, file_path: str, filename: str):
    text = ""
    
    try:
        if not os.path.exists(file_path):
            return 0, "", "File not found"
        
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif file_ext == '.docx':
            text = extract_text_from_docx(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            text = extract_text_from_pptx(file_path)
        elif file_ext == '.txt':
            text = extract_text_from_txt(file_path)
        elif file_ext in ['.md', '.markdown']:
            text = extract_text_from_md(file_path)
        elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff']:
            text = extract_text_from_image(file_path)
        elif file_ext == '.csv':
            text = extract_text_from_csv(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            text = extract_text_from_excel(file_path)
        elif file_ext in ['.zip', '.rar']:
            text = extract_text_from_zip(file_path)
        else:
            text = extract_text_from_txt(file_path)
        
        if text and len(text.strip()) > 10:
            sentences = re.split(r'(?<=[.!?])\s+', text)
            chunks = []
            current_chunk = ""
            
            for sentence in sentences:
                if len(current_chunk) + len(sentence) < 500:
                    current_chunk += sentence + " "
                else:
                    if current_chunk.strip():
                        chunks.append(current_chunk.strip())
                    current_chunk = sentence + " "
            
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            
            file_hash = hashlib.md5((filename + text[:100]).encode()).hexdigest()[:8]
            
            for i, chunk in enumerate(chunks):
                add_to_memory(conversation_id, chunk, filename, file_hash, i+1)
            
            return len(chunks), file_hash, text
        else:
            return 0, "", text if text else "Could not extract text"
            
    except Exception as e:
        return 0, "", f"Error: {str(e)}"

def retrieve_memory(conversation_id: str, query: str, k: int = 3):
    try:
        index, documents, metadata = load_vector_store(conversation_id)
        
        if index.ntotal == 0:
            return []
        
        query_embedding = get_embedding(query).reshape(1, -1)
        k = min(k, index.ntotal)
        distances, indices = index.search(query_embedding, k)
        
        results = []
        for i, idx in enumerate(indices[0]):
            if idx >= 0 and idx < len(documents):
                results.append({
                    "text": documents[idx],
                    "metadata": metadata[idx],
                    "distance": float(distances[0][i])
                })
        
        return results
    except:
        return []

def clear_memory(conversation_id: str = None):
    if conversation_id:
        try:
            INDEX_FILE, DATA_FILE, METADATA_FILE = get_conversation_store_path(conversation_id)
            index = faiss.IndexFlatL2(EMBEDDING_DIM)
            documents = []
            metadata = []
            
            faiss.write_index(index, INDEX_FILE)
            with open(DATA_FILE, "wb") as f:
                pickle.dump(documents, f)
            with open(METADATA_FILE, "wb") as f:
                pickle.dump(metadata, f)
        except:
            pass
    else:
        import shutil
        try:
            shutil.rmtree(VECTOR_DIR)
            os.makedirs(VECTOR_DIR, exist_ok=True)
        except:
            pass